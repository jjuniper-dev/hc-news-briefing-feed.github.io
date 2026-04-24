#!/usr/bin/env python3
import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
SPEC_PATH = ROOT / "diagrams" / "diagram-spec-v1.json"
VIEWS_DIR = ROOT / "diagrams" / "views"
OUTPUT_PATH = ROOT / "artifacts" / "ea-diagrams-v1.pptx"


def hex_rgb(value: str) -> RGBColor:
    value = value.lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def add_box(slide, x, y, w, h, text, fill, text_color, bold=False, size=14):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = fill

    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = text_color
    return shape


def main() -> None:
    spec = json.loads(SPEC_PATH.read_text())
    palette = spec["palette"]

    structure = hex_rgb(palette["structure"])
    components = hex_rgb(palette["components"])
    accent = hex_rgb(palette["accent"])
    text_color = hex_rgb(palette["text"])

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    view_paths = sorted(VIEWS_DIR.glob("view-*.json"))
    if len(view_paths) != 8:
        raise ValueError("Expected exactly 8 view manifests.")

    margin_x = Inches(0.6)
    title_y = Inches(0.25)
    col_top = Inches(1.2)
    col_bottom = Inches(5.2)
    col_gap = Inches(0.25)

    usable_w = prs.slide_width - (margin_x * 2)
    col_w = (usable_w - (col_gap * 3)) / 4

    bottom_y = Inches(5.75)
    bottom_h = Inches(1.2)

    for path in view_paths:
        manifest = json.loads(path.read_text())
        columns = manifest["columns"]
        if len(columns) != 4:
            raise ValueError(f"{path.name} must have 4 columns")

        for items in columns:
            if len(items) > 4:
                raise ValueError(f"{path.name} exceeds 4 items in a column")

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = hex_rgb(palette["background"])

        title = slide.shapes.add_textbox(margin_x, title_y, Inches(9), Inches(0.5))
        tframe = title.text_frame
        tframe.clear()
        p = tframe.paragraphs[0]
        p.text = manifest["title"]
        p.alignment = PP_ALIGN.LEFT
        p.runs[0].font.size = Pt(28)
        p.runs[0].font.bold = True
        p.runs[0].font.color.rgb = text_color

        col_centers = []
        for c_idx, items in enumerate(columns):
            x = margin_x + c_idx * (col_w + col_gap)
            col_centers.append(x + col_w / 2)
            rows = max(len(items), 1)
            item_h = Inches(0.65)
            total_h = rows * item_h + (rows - 1) * Inches(0.2)
            start_y = col_top + (col_bottom - col_top - total_h) / 2

            for i_idx, label in enumerate(items):
                y = start_y + i_idx * (item_h + Inches(0.2))
                fill = components
                if label == "Secure Gateway 2":
                    fill = structure
                elif label == "Microsoft Entra ID":
                    fill = accent
                add_box(slide, x, y, col_w, item_h, label, fill, text_color, bold=(fill != components), size=12)

        for i in range(3):
            line = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                int(col_centers[i] + col_w / 2 - Inches(0.1)),
                int((col_top + col_bottom) / 2),
                int(col_centers[i + 1] - col_w / 2 + Inches(0.1)),
                int((col_top + col_bottom) / 2),
            )
            line.line.color.rgb = structure
            line.line.width = Pt(2)

        bottom = add_box(
            slide,
            margin_x,
            bottom_y,
            usable_w,
            bottom_h,
            "Enterprise Data Platform",
            structure,
            RGBColor(255, 255, 255),
            bold=True,
            size=14,
        )

        items = ["Data Pipelines", "Data Products", "Analytics"]
        pill_w = Inches(2.0)
        pill_h = Inches(0.38)
        gap = Inches(0.2)
        total_w = len(items) * pill_w + (len(items) - 1) * gap
        start_x = margin_x + (usable_w - total_w) / 2
        for idx, lbl in enumerate(items):
            px = start_x + idx * (pill_w + gap)
            py = bottom_y + Inches(0.65)
            add_box(slide, px, py, pill_w, pill_h, lbl, components, text_color, size=10)

        overlays = manifest.get("overlays", [])
        if overlays:
            ov_w = Inches(2.15)
            ov_h = Inches(0.34)
            ov_x = prs.slide_width - margin_x - ov_w
            ov_y = Inches(0.3)
            for ov in overlays:
                add_box(slide, ov_x, ov_y, ov_w, ov_h, ov, RGBColor(255, 255, 255), structure, bold=True, size=10)
                ov_y += Inches(0.4)

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT_PATH)
    print(f"Generated {OUTPUT_PATH}")


if __name__ == "__main__":
    main()
